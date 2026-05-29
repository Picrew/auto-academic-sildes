from __future__ import annotations

from dataclasses import dataclass
from html import unescape
import json
from pathlib import Path
import re
import shutil
import subprocess
from tempfile import TemporaryDirectory

from pptx import Presentation
from pptx.util import Inches

from .ir import Deck
from .preview import contact_sheet
from .render_html import render_html


W, H = 13.333333, 7.5


TEXT_HARD_GATES = (
    "element-overlap",
    "text-block-overlap",
    "text-block-clearance-tight",
    "text-line-overlap",
    "text-line-height-tight",
    "text-self-overlap",
    "title-wrap-too-deep",
    "subtitle-wrap-too-deep",
    "text-clearance-tight",
    "text-image-overlap",
    "text-image-clearance-tight",
    "proof-notes-image-overlap",
    "artifact-notes-image-overlap",
    "notes-image-clearance-tight",
    "figure-caption-overlap",
    "figure-caption-clearance-tight",
    "text-clipped",
    "text-overflow",
    "container-overflow",
    "viewport-overflow",
    "pseudo-overlay-front",
)
IMAGE_HARD_GATES = (
    "missing-proof",
    "missing-proof-image",
    "image-not-loaded",
    "image-overflow",
    "image-contract-missing",
    "image-crop-missing",
    "image-caption-source-missing",
    "proof-too-small",
    "proof-image-too-small",
    "proof-image-rendered-too-small",
    "proof-image-upscaled-too-much",
    "proof-image-letterboxed-severe",
    "artifact-too-small",
    "artifact-image-too-small",
    "artifact-image-rendered-too-small",
    "artifact-image-upscaled-too-much",
    "artifact-image-letterboxed-severe",
    "image-object-fit-unsupported",
    "untyped-image",
    "untyped-vector-image",
    "untyped-background-image",
    "callout-outside-image",
    "callout-outside-source-image",
)
HARD_GATE_TYPES = TEXT_HARD_GATES + IMAGE_HARD_GATES
STRICT_WARNING_TYPES = (
    "proof-small",
    "proof-image-small",
    "proof-image-rendered-small",
    "proof-image-letterboxed",
    "artifact-small",
    "artifact-image-small",
    "artifact-image-rendered-small",
    "artifact-image-letterboxed",
    "callout-overlap",
    "content-underfilled",
    "decorative-image-too-small",
    "cover-image-letterboxed",
    "proof-caption-missing",
    "proof-caption-generic",
    "artifact-caption-generic",
    "caption-text-too-small",
    "caption-contrast-low",
    "label-contrast-low",
    "useful-fill-low",
    "artifact-role-underfeatured",
    "cadence-low-variety",
    "cadence-no-proof-surface",
    "title-wrap-deep",
    "subtitle-wrap-deep",
)


def effective_issue_level(issue: dict) -> str:
    kind = str(issue.get("type", ""))
    if kind in HARD_GATE_TYPES:
        return "error"
    return str(issue.get("level", "warn"))


@dataclass(frozen=True)
class HtmlPptxExport:
    html_path: Path
    images: tuple[Path, ...]
    pptx_path: Path
    contact_sheet_path: Path
    layout_report_path: Path
    layout_errors: int = 0
    strict_layout_warnings: int = 0


def find_chrome(explicit: str | None = None) -> str | None:
    candidates = [
        explicit,
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        shutil.which("google-chrome"),
        shutil.which("chromium"),
        shutil.which("chromium-browser"),
        shutil.which("microsoft-edge"),
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None


def screenshot_html_slides(
    html_path: Path,
    out_dir: Path,
    slide_count: int,
    width: int = 1920,
    height: int = 1080,
    chrome_path: str | None = None,
) -> tuple[Path, ...]:
    chrome = find_chrome(chrome_path)
    if not chrome:
        raise RuntimeError("Chrome/Chromium was not found. Set --chrome or install a Chromium-compatible browser.")
    out_dir.mkdir(parents=True, exist_ok=True)
    images: list[Path] = []
    for idx in range(1, slide_count + 1):
        target = out_dir / f"html-{idx:02d}.png"
        url = f"{html_path.resolve().as_uri()}?slide={idx}"
        command = [
            chrome,
            "--headless=new",
            "--disable-gpu",
            "--disable-background-networking",
            "--disable-extensions",
            "--disable-dev-shm-usage",
            "--no-first-run",
            "--hide-scrollbars",
            "--run-all-compositor-stages-before-draw",
            "--virtual-time-budget=1800",
            f"--window-size={width},{height}",
            f"--screenshot={target}",
            url,
        ]
        last_error: subprocess.TimeoutExpired | subprocess.CalledProcessError | None = None
        for timeout in (25, 35):
            try:
                subprocess.run(
                    command,
                    check=True,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    text=True,
                    timeout=timeout,
                )
                last_error = None
                break
            except (subprocess.TimeoutExpired, subprocess.CalledProcessError) as exc:
                last_error = exc
        if last_error:
            raise last_error
        images.append(target)
    return tuple(images)


def audit_html_layout(
    html_path: Path,
    slide_count: int,
    width: int = 1920,
    height: int = 1080,
    chrome_path: str | None = None,
    timeout_seconds: int = 25,
) -> tuple[dict, ...]:
    chrome = find_chrome(chrome_path)
    if not chrome:
        return tuple({"slide": idx, "issues": [{"level": "error", "type": "audit-unavailable", "target": "chrome", "detail": "Chrome/Chromium was not found."}]} for idx in range(1, slide_count + 1))
    audits: list[dict] = []
    pattern = re.compile(r'<pre id="layout-audit"[^>]*>(.*?)</pre>', re.DOTALL)
    with TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        for idx in range(1, slide_count + 1):
            url = f"{html_path.resolve().as_uri()}?slide={idx}&audit=1&expected_width={width}&expected_height={height}"
            audit_payload: str | None = None
            last_detail = "No audit block found."
            timed_out = False
            attempts: list[tuple[list[str], int]] = []
            for attempt in range(3):
                # Chrome's dump-dom path can ignore the requested viewport unless a screenshot is also requested.
                dummy_shot = tmp_dir / f"audit-{idx:02d}-{attempt + 1}.png"
                attempts.append(
                    (
                        [
                            chrome,
                            "--headless=new",
                            "--disable-gpu",
                            "--disable-background-networking",
                            "--disable-extensions",
                            "--disable-dev-shm-usage",
                            "--no-first-run",
                            "--hide-scrollbars",
                            f"--window-size={width},{height}",
                            f"--screenshot={dummy_shot}",
                            "--virtual-time-budget=5000",
                            "--dump-dom",
                            url,
                        ],
                        min(timeout_seconds + attempt * 10, 45),
                    )
                )
            attempts.append(
                (
                        [
                            chrome,
                            "--headless=new",
                            "--disable-gpu",
                            "--disable-background-networking",
                            "--disable-extensions",
                            "--disable-dev-shm-usage",
                            "--no-first-run",
                            "--hide-scrollbars",
                            f"--window-size={width},{height}",
                            "--virtual-time-budget=6000",
                            "--dump-dom",
                            url,
                    ],
                    min(timeout_seconds + 10, 45),
                )
            )
            for command, timeout in attempts:
                try:
                    result = subprocess.run(
                        command,
                        check=False,
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE,
                        text=True,
                        timeout=timeout,
                    )
                    timed_out = False
                    match = pattern.search(result.stdout)
                    if match:
                        audit_payload = unescape(match.group(1))
                        break
                    last_detail = (result.stderr or result.stdout or "No audit block found.")[:220]
                except subprocess.TimeoutExpired:
                    timed_out = True
                    last_detail = f"Chrome did not return within {timeout}s."
            if audit_payload is None and timed_out:
                audits.append(
                    {
                        "slide": idx,
                        "issues": [
                            {
                                "level": "error",
                                "type": "audit-timeout",
                                "target": "chrome",
                                "detail": last_detail,
                            }
                        ],
                    }
                )
                continue
            if audit_payload is None:
                audits.append(
                    {
                        "slide": idx,
                        "issues": [
                            {
                                "level": "error",
                                "type": "audit-missing",
                                "target": "layout-audit",
                                "detail": last_detail,
                            }
                        ],
                    }
                )
                continue
            try:
                audits.append(json.loads(audit_payload))
            except json.JSONDecodeError as exc:
                audits.append(
                    {
                        "slide": idx,
                        "issues": [
                            {
                                "level": "error",
                                "type": "audit-json",
                                "target": "layout-audit",
                                "detail": str(exc),
                            }
                        ],
                    }
                )
    return tuple(audits)


def cadence_issues(audits: tuple[dict, ...]) -> tuple[dict, ...]:
    if len(audits) < 5:
        return ()
    sequence = [str(audit.get("composition") or "unknown") for audit in audits]
    distinct = {item for item in sequence if item and item != "unknown"}
    issues: list[dict] = []
    if len(distinct) < 3:
        issues.append(
            {
                "level": "warn",
                "type": "cadence-low-variety",
                "target": "deck",
                "detail": f"Only {len(distinct)} composition type(s): {', '.join(sequence)}.",
            }
        )
    run_start = 0
    for idx in range(1, len(sequence) + 1):
        at_end = idx == len(sequence)
        changed = not at_end and sequence[idx] != sequence[run_start]
        if not (at_end or changed):
            continue
        run_len = idx - run_start
        if run_len >= 3 and sequence[run_start] not in {"unknown", "proof-led", "proof-atlas-spread", "proof-stage", "proof-dossier", "proof-ledger"}:
            issues.append(
                {
                    "level": "error",
                    "type": "cadence-repetition",
                    "target": "deck",
                    "detail": f"Slides {run_start + 1}-{idx} repeat `{sequence[run_start]}` composition.",
                }
            )
        run_start = idx
    evidence_like = sum(
        1
        for item in sequence
        if item
        in {
            "proof-led",
            "proof-atlas-spread",
            "proof-stage",
            "proof-dossier",
            "proof-ledger",
            "proof-marginalia",
            "proof-gallery-split",
            "proof-showcase",
            "artifact-content",
            "artifact-dossier",
            "artifact-rail",
            "artifact-marginalia",
            "artifact-ledger",
            "artifact-showcase",
        }
    )
    if evidence_like == 0:
        issues.append(
            {
                "level": "warn",
                "type": "cadence-no-proof-surface",
                "target": "deck",
                "detail": "No proof-led or artifact-content slide detected; source evidence may be too decorative or absent.",
            }
        )
    return tuple(issues)


def strict_warning_count(audits: tuple[dict, ...], deck_issues: tuple[dict, ...] = ()) -> int:
    count = 0
    for audit in audits:
        for issue in audit.get("issues", []):
            if effective_issue_level(issue) == "error":
                continue
            if issue.get("type") in STRICT_WARNING_TYPES:
                count += 1
    for issue in deck_issues:
        if effective_issue_level(issue) == "error":
            continue
        if issue.get("type") in STRICT_WARNING_TYPES:
            count += 1
    return count


def write_layout_audit_report(out_path: Path, audits: tuple[dict, ...]) -> tuple[Path, int]:
    errors = 0
    warnings = 0
    fills: list[float] = []
    useful_fills: list[float] = []
    lines = ["# HTML Layout Audit", ""]
    deck_issues = cadence_issues(audits)
    strict_warnings = strict_warning_count(audits, deck_issues)
    image_contract_rows: list[tuple[int | str, dict]] = []
    for audit in audits:
        fill = audit.get("metrics", {}).get("content_fill")
        if isinstance(fill, int | float):
            fills.append(float(fill))
        useful_fill = audit.get("metrics", {}).get("useful_fill")
        if isinstance(useful_fill, int | float):
            useful_fills.append(float(useful_fill))
        for image_contract in audit.get("images", []):
            if isinstance(image_contract, dict):
                image_contract_rows.append((audit.get("slide", "?"), image_contract))
        for issue in audit.get("issues", []):
            level = effective_issue_level(issue)
            if level == "error":
                errors += 1
            elif level == "warn":
                warnings += 1
    for issue in deck_issues:
        level = effective_issue_level(issue)
        if level == "error":
            errors += 1
        elif level == "warn":
            warnings += 1
    lines.extend(
        [
            f"- Slides audited: {len(audits)}",
            f"- Blocking errors: {errors}",
            f"- Warnings: {warnings}",
            f"- Strict quality warnings: {strict_warnings}",
            f"- Average content fill: {round(100 * sum(fills) / len(fills), 1) if fills else 'n/a'}%",
            f"- Average useful fill: {round(100 * sum(useful_fills) / len(useful_fills), 1) if useful_fills else 'n/a'}%",
            "",
        ]
    )
    if audits:
        sequence = [str(audit.get("composition") or "unknown") for audit in audits]
        distinct = sorted({item for item in sequence if item and item != "unknown"})
        lines.extend(
            [
                "## Deck Cadence",
                f"- Composition sequence: `{' -> '.join(sequence)}`",
                f"- Distinct cadences: {len(distinct)} ({', '.join(distinct) if distinct else 'none'})",
                "",
            ]
        )
        lines.extend(
            [
                "## Slide Metrics",
                "",
                "| Slide | Composition | Content fill | Useful fill | Issues |",
                "|---:|---|---:|---:|---:|",
            ]
        )
        for audit in audits:
            metrics = audit.get("metrics", {})
            fill = metrics.get("content_fill")
            useful_fill = metrics.get("useful_fill")
            fill_text = f"{round(float(fill) * 100, 1)}%" if isinstance(fill, int | float) else "n/a"
            useful_text = f"{round(float(useful_fill) * 100, 1)}%" if isinstance(useful_fill, int | float) else "n/a"
            issues_count = len(audit.get("issues", []))
            lines.append(
                f"| {audit.get('slide', '?')} | `{audit.get('composition') or 'unknown'}` | "
                f"{fill_text} | {useful_text} | {issues_count} |"
            )
        lines.append("")
    if image_contract_rows:
        lines.extend(
            [
                "## Image Contracts",
                "",
                "| Slide | Role | Rendered source | Visible area | Slot use | Loaded | Crop | Caption | Source | Target |",
                "|---:|---|---:|---:|---:|---|---|---|---|---|",
            ]
        )
        for slide_number, contract in image_contract_rows:
            rendered = f"{contract.get('rendered_width', 0)}x{contract.get('rendered_height', 0)}"
            visible = contract.get("visible_area")
            slot_use = contract.get("slot_use")
            visible_text = f"{round(float(visible) * 100, 1)}%" if isinstance(visible, int | float) else "n/a"
            slot_text = f"{round(float(slot_use) * 100, 1)}%" if isinstance(slot_use, int | float) else "n/a"
            crop = "yes" if contract.get("has_crop") else "no"
            caption = "yes" if contract.get("has_caption") else "no"
            source = "yes" if contract.get("has_source") else "no"
            loaded = "yes" if contract.get("loaded", True) else "no"
            role = contract.get("role", "image")
            target = contract.get("target", "img")
            lines.append(
                f"| {slide_number} | `{role}` | {rendered} | {visible_text} | {slot_text} | "
                f"{loaded} | {crop} | {caption} | {source} | `{target}` |"
            )
        lines.append("")
    if errors == 0 and warnings == 0:
        lines.append("No layout overlap, overflow, or proof-scale issues detected by the DOM audit.")
    else:
        lines.append("## Issues")
        if deck_issues:
            lines.append("### Deck cadence")
            for issue in deck_issues:
                level = effective_issue_level(issue)
                kind = issue.get("type", "cadence")
                target = issue.get("target", "deck")
                detail = issue.get("detail", "")
                lines.append(f"- `{level}` `{kind}` on `{target}`: {detail}")
        for audit in audits:
            issues = audit.get("issues", [])
            if not issues:
                continue
            metrics = audit.get("metrics", {})
            fill = metrics.get("content_fill")
            useful_fill = metrics.get("useful_fill")
            fill_parts = []
            if isinstance(fill, int | float):
                fill_parts.append(f"fill {round(float(fill) * 100)}%")
            if isinstance(useful_fill, int | float):
                fill_parts.append(f"useful {round(float(useful_fill) * 100)}%")
            fill_text = f" / {' / '.join(fill_parts)}" if fill_parts else ""
            lines.append(f"### Slide {audit.get('slide', '?')}{fill_text}")
            for issue in issues:
                level = effective_issue_level(issue)
                kind = issue.get("type", "layout")
                target = issue.get("target", "unknown")
                detail = issue.get("detail", "")
                lines.append(f"- `{level}` `{kind}` on `{target}`: {detail}")
    lines.extend(
        [
            "",
            "## Hard Gates",
            f"- Text gates: {', '.join(f'`{gate}`' for gate in TEXT_HARD_GATES)}.",
            f"- Image gates: {', '.join(f'`{gate}`' for gate in IMAGE_HARD_GATES)}.",
            f"- Strict warning gates: {', '.join(f'`{gate}`' for gate in STRICT_WARNING_TYPES)}.",
            "- Any text gate must be fixed before treating the deck as presentable.",
            "- Any `text-line-height-tight` error means wrapped text is too compressed even if glyphs have not visibly collided yet; increase line-height, lower type scale, or split the slide.",
            "- Any `title-wrap-too-deep` or `subtitle-wrap-too-deep` error means the copy is using the visual canvas as a document; shorten it or split the slide before CSS polish.",
            "- Any `text-image-overlap` or `text-image-clearance-tight` error means the text and visual proof are fighting for the same space; shorten the copy, crop tighter, or switch composition.",
            "- Any image gate means the visual evidence is not usable yet.",
            "- Every proof/artifact image should appear in `Image Contracts` with a role, rendered source size, visible area, crop, caption, and source status.",
            "- In strict export, strict warning gates also fail `--fail-on-layout`; they indicate an inserted image, callout, or page density is not ready for delivery.",
            "- Any `proof-image-letterboxed` or `artifact-image-letterboxed` warning means the crop or slot shape should be reviewed before judging the design tasteful.",
            "- Any severe letterboxing error means the image was inserted but the readable source pixels are too small for slide-scale proof.",
            "- Any `proof-image-small` warning means the proof surface is present but underfeatured; switch to a larger proof layout or crop tighter before final polish.",
            "- Any proof/artifact caption warning means the image was inserted without enough source context or readable caption styling; add a specific claim caption, compact source line, and sufficient caption contrast/size.",
            "- Evidence-like slides need a proof surface large enough to inspect; a tiny screenshot is a failed proof, not a style choice.",
            "- Content slides that carry screenshots should use an artifact panel large enough to read; otherwise remove the image or redesign the slide.",
            "- Any `useful-fill-low` warning means the slide is visually occupied by frames or whitespace but lacks enough useful text/source evidence.",
            "- Any `artifact-role-underfeatured` warning means a source-heavy artifact such as a repo, UI, table, paper, benchmark, or project page is too small for inspection.",
            "- Decks of five or more slides should show at least three composition cadences and avoid three repeated non-proof compositions in a row.",
            "- This DOM audit is conservative; manual contact-sheet review still checks taste, crop quality, and whether the image proves the claim.",
        ]
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(lines), encoding="utf-8")
    blocking_errors = 0
    for audit in audits:
        for issue in audit.get("issues", []):
            if effective_issue_level(issue) != "error" or issue.get("type") == "content-underfilled":
                continue
            blocking_errors += 1
    for issue in deck_issues:
        if effective_issue_level(issue) == "error":
            blocking_errors += 1
    return out_path, blocking_errors


def write_skipped_layout_audit_report(out_path: Path) -> tuple[Path, int]:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        "\n".join(
            [
                "# HTML Layout Audit",
                "",
                "- Status: skipped for fast visual exploration.",
                "",
                "Run `uv run academic-deck html-pptx ... --fail-on-layout` without `--skip-layout-audit` on shortlisted variants.",
            ]
        ),
        encoding="utf-8",
    )
    return out_path, 0


def create_image_pptx(images: tuple[Path, ...], output_path: Path) -> Path:
    if not images:
        raise ValueError("No images supplied for PPTX export.")
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def export_html_image_pptx(
    deck: Deck,
    output_dir: Path,
    width: int = 1920,
    height: int = 1080,
    chrome_path: str | None = None,
    skip_layout_audit: bool = False,
) -> HtmlPptxExport:
    output_dir.mkdir(parents=True, exist_ok=True)
    html_path = render_html(output_dir / "html", deck)
    image_dir = output_dir / "html-shots"
    images = screenshot_html_slides(html_path, image_dir, len(deck.slides), width=width, height=height, chrome_path=chrome_path)
    pptx_path = create_image_pptx(images, output_dir / f"{deck.deck_id}-html-image.pptx")
    contact = contact_sheet(list(images), output_dir / "html-contact-sheet.png")
    if skip_layout_audit:
        layout_report, layout_errors = write_skipped_layout_audit_report(output_dir / "layout-audit-report.md")
        strict_layout_warnings = 0
    else:
        audits = audit_html_layout(html_path, len(deck.slides), width=width, height=height, chrome_path=chrome_path)
        layout_report, layout_errors = write_layout_audit_report(output_dir / "layout-audit-report.md", audits)
        strict_layout_warnings = strict_warning_count(audits, cadence_issues(audits))
    return HtmlPptxExport(
        html_path=html_path,
        images=images,
        pptx_path=pptx_path,
        contact_sheet_path=contact,
        layout_report_path=layout_report,
        layout_errors=layout_errors,
        strict_layout_warnings=strict_layout_warnings,
    )
