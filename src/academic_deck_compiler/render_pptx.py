from __future__ import annotations

from hashlib import sha1
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .assets import crop_asset, resolve_asset
from .content import Slide
from .ir import Deck, DEFAULT_DECK
from .patterns import parse_matrix_rows
from .style import THEME


W, H = 13.333333, 7.5
M = 0.58
ACTIVE_DECK: Deck = DEFAULT_DECK
ACTIVE_CACHE_DIR: Path | None = None


def asset_path(image_name: str) -> Path:
    return resolve_asset(ACTIVE_DECK, image_name)


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def emu(value: float):
    return Inches(value)


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color: str, width: float = 0.7, transparency: int = 0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def no_line(shape) -> None:
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill=None, line=None, radius=False, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, emu(x), emu(y), emu(w), emu(h))
    try:
        shape.shadow.inherit = False
    except AttributeError:
        pass
    if fill:
        set_fill(shape, fill, transparency)
    else:
        shape.fill.background()
    if line:
        set_line(shape, line)
    else:
        no_line(shape)
    return shape


def add_oval(slide, x, y, w, h, fill=None, line=None, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(x), emu(y), emu(w), emu(h))
    try:
        shape.shadow.inherit = False
    except AttributeError:
        pass
    if fill:
        set_fill(shape, fill, transparency)
    else:
        shape.fill.background()
    if line:
        set_line(shape, line)
    else:
        no_line(shape)
    return shape


def add_line(slide, x1, y1, x2, y2, color=None, width=0.7, transparency: int = 0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    set_line(line, color or THEME.hair, width, transparency)
    return line


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    color: str | None = None,
    font: str | None = None,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font or THEME.sans
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color or THEME.ink)
    return box


def add_multiline(slide, lines, x, y, w, h, size=16, color=None, bullet=False, gap=0.12):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    for i, line in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.space_after = Pt(gap * 24)
        p.line_spacing = 1.06
        if bullet:
            p.level = 0
            p.text = ""
            marker = p.add_run()
            marker.text = "-> "
            marker.font.name = THEME.mono
            marker.font.size = Pt(max(10, size - 2))
            marker.font.bold = True
            marker.font.color.rgb = rgb(THEME.copper)
            run = p.add_run()
        else:
            run = p.add_run()
        run.text = line
        run.font.name = THEME.sans
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color or THEME.ink)
    return box


def add_kicker(slide, text: str, color: str | None = None) -> None:
    add_text(slide, text.upper(), M, 0.42, 4.8, 0.22, size=9, color=color or THEME.copper, font=THEME.mono, bold=True)


def add_title(slide, title: str, y=0.66, h=0.68, size=23, width=10.5) -> None:
    add_text(slide, title, M, y, width, h, size=size, color=THEME.ink, font=THEME.sans, bold=True, line_spacing=0.96)
    add_line(slide, M, y + h + 0.12, W - M, y + h + 0.12, THEME.hair, 0.6)


def add_footer(slide, idx: int, total: int) -> None:
    footer = ACTIVE_DECK.footer or ACTIVE_DECK.title
    add_text(slide, footer, M, 7.1, 4.6, 0.16, size=8, color=THEME.muted, font=THEME.sans)
    add_text(slide, f"{idx}/{total}", W - 0.98, 7.1, 0.42, 0.16, size=8, color=THEME.muted, font=THEME.sans, align=PP_ALIGN.RIGHT)


def add_image(slide, image_name: str, x, y, w, h, border=True):
    path = asset_path(image_name)
    if not path.exists():
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    target_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= target_ratio:
        pic_h = h
        pic_w = h * img_ratio
        px = x - (pic_w - w) / 2
        py = y
    else:
        pic_w = w
        pic_h = w / img_ratio
        px = x
        py = y - (pic_h - h) / 2
    pic = slide.shapes.add_picture(str(path), emu(px), emu(py), width=emu(pic_w), height=emu(pic_h))
    try:
        pic.shadow.inherit = False
    except AttributeError:
        pass
    if border:
        add_rect(slide, x, y, w, h, fill=None, line=THEME.hair)
    return pic


def add_image_fit(slide, image_name: str, x, y, w, h, border=True):
    path = asset_path(image_name)
    if not path.exists():
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    target_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= target_ratio:
        pic_w = w
        pic_h = w / img_ratio
        px = x
        py = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * img_ratio
        px = x + (w - pic_w) / 2
        py = y
    if border:
        add_rect(slide, x, y, w, h, fill="FFFFFF", line=THEME.hair)
    pic = slide.shapes.add_picture(str(path), emu(px), emu(py), width=emu(pic_w), height=emu(pic_h))
    try:
        pic.shadow.inherit = False
    except AttributeError:
        pass
    return pic


def evidence_image_path(slide_data: Slide) -> Path | None:
    evidence = slide_data.evidence
    image_name = evidence.image if evidence and evidence.image else slide_data.image
    if not image_name:
        return None
    source = asset_path(image_name)
    if not source.exists():
        return None
    if evidence and evidence.crop:
        key = sha1(f"{source}|{evidence.crop}".encode("utf-8")).hexdigest()[:16]
        return crop_asset(source, evidence.crop, ACTIVE_CACHE_DIR or source.parent / ".cache", key)
    return source


def add_image_path_fit(slide, path: Path, x, y, w, h, border=True):
    with Image.open(path) as im:
        iw, ih = im.size
    target_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= target_ratio:
        pic_w = w
        pic_h = w / img_ratio
        px = x
        py = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * img_ratio
        px = x + (w - pic_w) / 2
        py = y
    if border:
        add_rect(slide, x, y, w, h, fill="FFFFFF", line=THEME.hair)
    pic = slide.shapes.add_picture(str(path), emu(px), emu(py), width=emu(pic_w), height=emu(pic_h))
    try:
        pic.shadow.inherit = False
    except AttributeError:
        pass
    return pic


def add_callout(slide, text: str, anchor_x: float, anchor_y: float, box_x: float, box_y: float, w: float = 2.1):
    add_line(slide, anchor_x, anchor_y, box_x, box_y + 0.16, THEME.copper, 1.0)
    add_rect(slide, box_x, box_y, w, 0.44, fill="FFFFFF", line=THEME.copper)
    add_text(slide, text, box_x + 0.12, box_y + 0.12, w - 0.24, 0.12, size=8, color=THEME.ink, font=THEME.sans)


def add_pin(slide, x: float, y: float, label: str, dark: bool = False) -> None:
    fill = THEME.copper
    ink = THEME.graphite if dark else THEME.paper
    add_oval(slide, x - 0.105, y - 0.105, 0.21, 0.21, fill=fill, line=None)
    add_text(
        slide,
        label,
        x - 0.062,
        y - 0.055,
        0.125,
        0.09,
        size=6,
        color=ink,
        font=THEME.mono,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_note_ledger(slide, notes, x: float, y: float, w: float, dark: bool = False) -> None:
    text_color = "DDE5EA" if dark else THEME.ink
    muted = "8A96A8" if dark else THEME.muted
    line_color = "53606B" if dark else THEME.hair
    for i, note in enumerate(notes[:3]):
        yy = y + i * 0.54
        add_text(slide, f"{i + 1:02d}", x, yy + 0.01, 0.32, 0.1, size=7, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, note, x + 0.42, yy - 0.01, w - 0.42, 0.22, size=8, color=text_color, font=THEME.sans)
        add_line(slide, x, yy + 0.33, x + w, yy + 0.33, line_color, 0.45, transparency=18)
    if not notes:
        add_text(slide, "pin notes", x, y, w, 0.12, size=7, color=muted, font=THEME.mono, bold=True)


def add_evidence_pins(slide, evidence, img_x: float, img_y: float, img_w: float, img_h: float, dark: bool = False):
    notes = []
    if not evidence:
        return notes
    for i, callout in enumerate(evidence.callouts[:3]):
        add_pin(slide, img_x + callout.x * img_w, img_y + callout.y * img_h, str(i + 1), dark=dark)
        notes.append(callout.text)
    return notes


def add_source_note(slide, text: str, x: float, y: float, w: float):
    if text:
        add_text(slide, text, x, y, w, 0.12, size=7, color=THEME.muted, font=THEME.mono)


def add_node(slide, text, x, y, w, h, fill="FFFFFF", line=None, color=None, size=10, mono=False):
    add_rect(slide, x, y, w, h, fill=fill, line=line or THEME.hair)
    add_text(
        slide,
        text,
        x + 0.12,
        y + h / 2 - 0.055,
        w - 0.24,
        0.12,
        size=size,
        color=color or THEME.ink,
        font=THEME.mono if mono else THEME.sans,
        align=PP_ALIGN.CENTER,
    )


def add_value_label(slide, value, label, x, y, w=1.4, dark=False):
    size = 17 if len(value) > 6 else 19
    add_text(slide, value, x, y, w, 0.28, size=size, color=THEME.paper if dark else THEME.ink, font=THEME.serif)
    add_text(slide, label, x, y + 0.34, w, 0.18, size=7, color="AAB6BF" if dark else THEME.muted, font=THEME.sans)


def canvas(slide, dark=False) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(THEME.graphite if dark else THEME.paper)


def add_grid(slide, dark=False) -> None:
    if not dark:
        return
    color = "344260"
    x = 0.0
    while x <= W + 0.01:
        add_line(slide, x, 0, x, H, color, 0.35, transparency=58)
        x += 1.12
    y = 0.0
    while y <= H + 0.01:
        add_line(slide, 0, y, W, y, color, 0.35, transparency=58)
        y += 0.9


def render_cover(slide, s: Slide) -> None:
    canvas(slide, dark=True)
    add_rect(slide, 0, 0, W, H, fill=THEME.graphite)
    add_grid(slide, dark=True)
    add_line(slide, M, 0.62, W - M, 0.62, THEME.copper, 0.8)
    add_text(slide, s.kicker or "RESEARCH DECK", M, 0.92, 4.8, 0.2, size=9, color=THEME.copper, font=THEME.mono, bold=True)
    hero = s.title.replace(" need ", "\nneed ", 1).replace(" more than ", "\nmore than ", 1)
    add_text(slide, hero, M, 1.36, 8.15, 1.7, size=38, color=THEME.paper, font=THEME.serif, line_spacing=0.88)
    add_text(slide, ACTIVE_DECK.title, M, 3.58, 3.2, 0.32, size=19, color=THEME.copper, font=THEME.serif)
    add_text(slide, ACTIVE_DECK.subtitle, M, 3.98, 6.65, 0.24, size=11, color="DDE5EA", font=THEME.sans)
    add_text(slide, s.subtitle, M, 4.42, 6.55, 0.32, size=9, color="8A96A8", font=THEME.sans)
    image_top = 1.16
    focus_top = 1.48
    focus_h = 3.42
    if s.image:
        source = asset_path(str(s.image))
        if source.exists():
            add_rect(slide, 8.74, 1.08, 3.48, 2.02, fill="FFFFFF", line="53606B")
            add_image_path_fit(slide, source, 8.84, 1.18, 3.28, 1.82, border=False)
            add_text(slide, "source surface", 8.84, 3.22, 1.15, 0.1, size=7, color="8A96A8", font=THEME.mono, bold=True)
            image_top = 3.46
            focus_top = 3.7
            focus_h = 1.12
    add_rect(slide, 9.65, image_top, 2.55, focus_h, fill=None, line="53606B")
    add_text(slide, "focus", 9.9, focus_top, 0.7, 0.12, size=7, color="8A96A8", font=THEME.mono, bold=True)
    for i, label in enumerate(s.labels[:4]):
        y = focus_top + 0.4 + i * 0.24 if s.image else 1.88 + i * 0.48
        add_text(slide, f"{i+1:02d}", 9.9, y, 0.28, 0.1, size=7, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, label, 10.28, y - 0.02, 1.55, 0.12, size=7 if s.image else 8, color=THEME.paper, font=THEME.sans)
    add_line(slide, M, 5.32, W - M, 5.32, "53606B", 0.65)
    for i, (value, label) in enumerate(s.metrics[:3]):
        mx = M + i * 3.68
        add_text(slide, value, mx, 5.68, 2.25, 0.34, size=25, color=THEME.copper, font=THEME.serif)
        add_text(slide, label, mx, 6.12, 2.25, 0.18, size=8, color="8A96A8", font=THEME.sans)
    add_line(slide, M, 6.82, W - M, 6.82, "53606B", 0.65)


def render_thesis(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide)
    add_kicker(slide, s.kicker)
    add_title(slide, s.title, size=25, width=11.1)
    add_text(slide, s.subtitle, M, 1.92, 5.95, 0.84, size=16, color=THEME.ink, font=THEME.serif, line_spacing=0.96)
    labels = ("failure", "signal", "policy", "system")
    for i, label in enumerate(labels):
        y = 3.28 + i * 0.6
        add_text(slide, f"0{i+1}", M, y, 0.35, 0.11, size=8, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, label, 1.04, y - 0.06, 1.5, 0.16, size=13, color=THEME.ink, font=THEME.sans, bold=True)
        if i < len(labels) - 1:
            add_line(slide, 1.78, y + 0.06, 2.3, y + 0.06, THEME.copper, 0.8)
    for i, bullet in enumerate(s.bullets[:4]):
        name, desc = (bullet.split(":", 1) + [""])[:2] if ":" in bullet else (f"Thesis {i+1}", bullet)
        y = 2.14 + i * 0.82
        add_line(slide, 6.82, y - 0.08, 11.95, y - 0.08, THEME.hair, 0.55)
        add_text(slide, f"{i+1:02d}", 6.82, y + 0.08, 0.35, 0.11, size=8, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, name.strip(), 7.3, y, 1.35, 0.16, size=12, color=THEME.ink, font=THEME.sans, bold=True)
        add_text(slide, desc.strip()[:92], 8.62, y, 3.2, 0.26, size=9, color=THEME.muted, font=THEME.sans)
    add_text(slide, "Research taste becomes useful when it survives engineering pressure.", 6.82, 5.7, 4.82, 0.2, size=13, color=THEME.ink, font=THEME.serif)
    add_footer(slide, idx, total)


def render_map(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide, dark=True)
    add_grid(slide, dark=True)
    add_kicker(slide, s.kicker)
    add_text(slide, s.title, M, 0.72, 9.7, 0.58, size=27, color=THEME.paper, font=THEME.serif)
    add_text(slide, s.subtitle, M, 1.55, 5.9, 0.32, size=11, color="C7D1D8", font=THEME.sans)
    stage_names = tuple(s.labels[:5]) if len(s.labels) >= 5 else ("Question", "Failure", "Signal", "Intervention", "Result")
    stage_prompts = ("what broke?", "what can supervise it?", "what changes?", "what verifies it?", "what ships?")
    xs = tuple(0.9 + i * 2.28 for i in range(len(stage_names)))
    top_y = 3.08
    stages = tuple(zip(stage_names, stage_prompts))
    for i, ((stage, sub), x) in enumerate(zip(stages, xs)):
        add_text(slide, f"0{i+1}", x, top_y - 0.44, 0.4, 0.12, size=8, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, stage, x, top_y, 1.42, 0.25, size=19, color=THEME.paper, font=THEME.serif)
        add_text(slide, sub, x, top_y + 0.48, 1.45, 0.13, size=7, color="8A96A8", font=THEME.mono)
        if i < len(stages) - 1:
            add_line(slide, x + 1.22, top_y + 0.18, x + 2.0, top_y + 0.18, THEME.copper, 0.85)
    cases = []
    for bullet in s.bullets[:4]:
        if ":" in bullet:
            name, desc = bullet.split(":", 1)
        else:
            words = bullet.split()
            name, desc = " ".join(words[:2]), " ".join(words[2:8])
        cases.append((name[:22], desc.strip()[:34]))
    if not cases:
        cases = [(label, "control surface") for label in stage_names]
    for i, (name, desc) in enumerate(cases):
        x = 0.9 + i * 2.95
        add_line(slide, x, 5.22, x + 2.1, 5.22, "2E3D5C", 0.6)
        add_text(slide, name, x, 5.46, 1.65, 0.14, size=10, color=THEME.paper, font=THEME.sans, bold=True)
        add_text(slide, "control surface", x, 5.7, 1.35, 0.1, size=6, color="8A96A8", font=THEME.mono)
    add_footer(slide, idx, total)


def render_split(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide)
    add_kicker(slide, s.kicker)
    add_title(slide, s.title, size=25, width=10.8)
    add_text(slide, s.subtitle, M, 1.86, 5.8, 0.44, size=11, color=THEME.muted, font=THEME.sans)
    add_rect(slide, 0.9, 2.72, 4.85, 2.35, fill="FFFFFF", line=THEME.hair)
    left_label = s.labels[0] if len(s.labels) >= 1 else "retrieval as add-on"
    right_label = s.labels[1] if len(s.labels) >= 2 else "search as trained behavior"
    left_note = s.bullets[0] if s.bullets else "tool exists, but the policy can still learn to skip it"
    right_note = s.bullets[1] if len(s.bullets) > 1 else "reward and rollout design make evidence-first the stable path"
    add_text(slide, left_label, 1.15, 3.0, 2.2, 0.2, size=17, color=THEME.ink, font=THEME.serif)
    add_text(slide, left_note[:82], 1.16, 3.42, 3.7, 0.16, size=9, color=THEME.muted, font=THEME.mono)
    add_node(slide, "question", 1.2, 4.05, 1.15, 0.34, mono=True)
    add_line(slide, 2.42, 4.22, 3.02, 4.22, THEME.copper, 0.9)
    add_node(slide, "answer", 3.1, 4.05, 1.15, 0.34, fill=THEME.soft_copper, mono=True)
    add_text(slide, "shortcut attractor", 1.23, 4.72, 2.4, 0.16, size=8, color=THEME.copper, font=THEME.mono, bold=True)
    add_rect(slide, 6.22, 2.72, 5.8, 2.35, fill=THEME.graphite, line=None)
    add_text(slide, right_label, 6.55, 3.0, 2.9, 0.2, size=17, color=THEME.paper, font=THEME.serif)
    add_text(slide, right_note[:82], 6.55, 3.42, 4.2, 0.16, size=9, color="AAB6BF", font=THEME.mono)
    for i, node in enumerate(("question", "search", "evidence", "answer")):
        x = 6.65 + i * 1.22
        add_node(slide, node, x, 4.05, 0.92, 0.34, fill="202832", line="4A5967", color=THEME.paper, size=7, mono=True)
        if i < 3:
            add_line(slide, x + 0.96, 4.22, x + 1.17, 4.22, THEME.copper, 0.9)
    interventions = ("staged reward", "forced exploration", "token masking")
    for i, item in enumerate(interventions):
        add_text(slide, f"{i+1:02d}", 1.05 + i * 3.15, 5.72, 0.35, 0.12, size=8, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, item, 1.48 + i * 3.15, 5.68, 2.2, 0.18, size=12, color=THEME.ink, font=THEME.serif)
    add_text(slide, "The contribution is the policy geometry: what the model is rewarded to do before it speaks.", M, 6.42, 8.3, 0.18, size=9, color=THEME.muted, font=THEME.mono)
    add_footer(slide, idx, total)


def render_process(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide)
    add_kicker(slide, s.kicker)
    add_title(slide, s.title, size=25, width=10.7)
    add_text(slide, s.subtitle, M, 1.84, 5.4, 0.35, size=11, color=THEME.muted, font=THEME.sans)
    add_text(slide, "lever", 1.0, 2.52, 0.8, 0.1, size=7, color=THEME.muted, font=THEME.mono, bold=True)
    add_text(slide, "failure prevented", 4.2, 2.52, 1.4, 0.1, size=7, color=THEME.muted, font=THEME.mono, bold=True)
    add_text(slide, "training signal", 8.08, 2.52, 1.3, 0.1, size=7, color=THEME.muted, font=THEME.mono, bold=True)
    add_line(slide, 0.92, 2.82, 12.05, 2.82, THEME.hair, 0.55)
    for i, (label, bullet) in enumerate(zip(s.labels, s.bullets)):
        y = 3.16 + i * 0.92
        add_text(slide, f"0{i+1}", 1.0, y + 0.04, 0.35, 0.1, size=8, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, label, 1.45, y - 0.02, 2.2, 0.18, size=15, color=THEME.ink, font=THEME.sans, bold=True)
        presets = {
            "Reward shaping": ("sparse final-answer reward", "dense signal for format, grounding, correctness"),
            "Forced exploration": ("early no-search collapse", "search sampled before the policy locks in"),
            "Retrieved-token masking": ("tool text contaminates loss", "loss only on model-authored tokens"),
        }
        failure, signal = presets.get(label, ("shortcut learning", bullet[:56]))
        add_text(slide, failure, 4.2, y, 2.65, 0.18, size=10, color=THEME.muted, font=THEME.sans)
        add_text(slide, signal, 8.08, y, 3.2, 0.2, size=10, color=THEME.muted, font=THEME.sans)
        add_line(slide, 0.92, y + 0.5, 12.05, y + 0.5, THEME.hair, 0.55)
    add_text(slide, "The levers make search creditable, sampled, and cleanly supervised.", 1.0, 6.14, 5.8, 0.14, size=9, color=THEME.muted, font=THEME.mono)
    add_footer(slide, idx, total)


def render_evidence(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide)
    add_kicker(slide, s.kicker)
    add_title(slide, s.title, size=24, width=11.2)
    add_text(slide, s.subtitle, M, 1.58, 5.85, 0.42, size=10, color=THEME.muted, font=THEME.sans)
    names = tuple(s.labels[:3]) if len(s.labels) >= 3 else ("Claim", "Artifact", "Implication")
    descriptions = tuple(item.split(";")[0] for item in s.bullets[:3]) if len(s.bullets) >= 3 else (
        "State the judgment.",
        "Show the cropped proof object.",
        "Name the implication.",
    )
    if names == ("Taxonomy", "Checker", "Metrics"):
        descriptions = (
            "2,000 prompts / 7 corpora",
            "taxonomy plus evidence chains",
            "CED + GRR diagnostics",
        )
    for i, (name, desc) in enumerate(zip(names, descriptions)):
        y = 2.48 + i * 0.78
        add_text(slide, f"0{i+1}", 0.86, y, 0.35, 0.1, size=8, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, name, 1.28, y - 0.06, 1.45, 0.16, size=12, color=THEME.ink, font=THEME.sans, bold=True)
        add_text(slide, desc[:48], 1.28, y + 0.22, 2.75, 0.18, size=8, color=THEME.muted, font=THEME.sans)
        add_line(slide, 0.86, y + 0.52, 4.12, y + 0.52, THEME.hair, 0.5)
    image_path = evidence_image_path(s)
    if image_path:
        with Image.open(image_path) as im:
            image_ratio = im.size[0] / im.size[1]
        if image_ratio < 1.35:
            img_x, img_y, img_w, img_h = 5.18, 1.94, 4.65, 4.42
            note_x, note_y, note_w = 10.25, 2.12, 2.1
        else:
            img_x, img_y, img_w, img_h = 4.14, 1.78, 7.82, 4.28
            note_x, note_y, note_w = 0.86, 5.2, 3.1
        add_rect(slide, img_x - 0.14, img_y - 0.14, img_w + 0.28, img_h + 0.72, fill="FFFFFF", line=THEME.hair)
        add_image_path_fit(slide, image_path, img_x, img_y, img_w, img_h, border=False)
        evidence = s.evidence
        caption = evidence.caption if evidence else ""
        source = evidence.source if evidence else ""
        add_text(slide, caption or "Evidence artifact cropped for inspection", img_x, img_y + img_h + 0.2, img_w, 0.16, size=8, color=THEME.ink, font=THEME.sans, bold=True)
        add_source_note(slide, source, img_x, img_y + img_h + 0.46, img_w)
        notes = add_evidence_pins(slide, evidence, img_x, img_y, img_w, img_h, dark=False)
        add_text(slide, "what to inspect", note_x, note_y - 0.34, note_w, 0.1, size=7, color=THEME.muted, font=THEME.mono, bold=True)
        add_note_ledger(slide, notes, note_x, note_y, note_w, dark=False)
    else:
        img_x, img_y, img_w, img_h = 4.8, 2.02, 6.7, 3.86
        add_rect(slide, img_x - 0.14, img_y - 0.14, img_w + 0.28, img_h + 0.28, fill="FFFFFF", line=THEME.hair)
        add_text(slide, "drop in one proof artifact", 5.52, 3.2, 3.0, 0.22, size=17, color=THEME.ink, font=THEME.sans, bold=True)
        add_text(slide, "figure / table / trace / UI crop", 5.54, 3.7, 2.2, 0.12, size=8, color=THEME.copper, font=THEME.mono, bold=True)
        add_line(slide, 5.54, 4.24, 10.78, 4.24, THEME.hair, 0.7)
        add_text(slide, "The renderer leaves this space quiet until a real asset earns the visual weight.", 5.54, 4.44, 4.1, 0.16, size=8, color=THEME.muted, font=THEME.sans)
    for i, (value, label) in enumerate(s.metrics):
        add_value_label(slide, value, label, 0.86 + i * 1.42, 5.76, w=1.28)
    add_footer(slide, idx, total)


def render_matrix(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide, dark=True)
    add_grid(slide, dark=True)
    add_kicker(slide, s.kicker, color=THEME.copper)
    add_text(slide, s.title, M, 0.68, 9.9, 0.62, size=25, color=THEME.paper, font=THEME.serif, line_spacing=0.92)
    add_line(slide, M, 1.7, W - M, 1.7, "53606B", 0.65)
    add_text(slide, s.subtitle, M, 2.0, 4.85, 0.96, size=11, color="C7D1D8", font=THEME.sans)
    if s.metrics:
        add_text(slide, s.metrics[1][0] if len(s.metrics) > 1 else s.metrics[0][0], M, 3.34, 2.4, 0.42, size=27, color=THEME.copper, font=THEME.serif)
        add_text(slide, s.metrics[1][1] if len(s.metrics) > 1 else s.metrics[0][1], M, 3.86, 2.25, 0.16, size=8, color="8A96A8", font=THEME.sans)
        add_text(slide, s.metrics[0][0], M, 4.62, 1.5, 0.28, size=19, color=THEME.paper, font=THEME.serif)
        add_text(slide, s.metrics[0][1], M, 4.94, 2.0, 0.14, size=8, color="8A96A8", font=THEME.sans)
    x0, y0 = 5.15, 2.08
    add_text(slide, "surface", x0, y0, 1.1, 0.1, size=7, color="8A96A8", font=THEME.mono, bold=True)
    add_text(slide, "failure controlled", x0 + 2.55, y0, 1.5, 0.1, size=7, color="8A96A8", font=THEME.mono, bold=True)
    add_text(slide, "artifact", x0 + 5.15, y0, 1.0, 0.1, size=7, color="8A96A8", font=THEME.mono, bold=True)
    add_line(slide, x0, y0 + 0.28, 12.08, y0 + 0.28, "53606B", 0.55)
    rows = parse_matrix_rows(s.labels[:7], s.bullets)
    for i, (label, failure, artifact) in enumerate(rows):
        y = y0 + 0.58 + i * 0.48
        add_text(slide, f"{i+1:02d}", x0, y, 0.35, 0.1, size=7, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, label, x0 + 0.46, y - 0.02, 1.85, 0.12, size=10, color=THEME.paper, font=THEME.sans, bold=True)
        add_text(slide, failure[:38], x0 + 2.55, y - 0.01, 2.25, 0.1, size=8, color="C7D1D8", font=THEME.sans)
        add_text(slide, artifact[:30], x0 + 5.15, y - 0.01, 1.8, 0.1, size=7, color="8A96A8", font=THEME.mono)
        add_line(slide, x0, y + 0.24, 12.08, y + 0.24, "2E3D5C", 0.45)
    add_text(slide, f"{idx}/{total}", W - 0.98, 7.1, 0.42, 0.16, size=8, color="AAB6BF", font=THEME.sans, align=PP_ALIGN.RIGHT)


def render_system(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide)
    add_kicker(slide, s.kicker)
    add_title(slide, s.title, size=26, width=10.8)
    add_text(slide, s.subtitle, M, 1.82, 5.3, 0.58, size=12, color=THEME.muted, font=THEME.sans)
    add_multiline(slide, s.bullets, M, 2.76, 4.6, 1.68, size=11, bullet=True)
    add_rect(slide, 5.82, 2.18, 5.95, 2.82, fill="FFFFFF", line=THEME.hair)
    lanes = (("text", 6.15, 2.55), ("face position", 7.65, 2.55), ("surface", 9.34, 2.55), ("topology", 10.38, 3.55), ("valid CAD", 8.32, 4.22))
    for label, x, y in lanes:
        fill = THEME.soft_blue if label in ("face position", "valid CAD") else "FFFFFF"
        add_node(slide, label, x, y, 1.1, 0.32, fill=fill, mono=True, size=7)
    add_line(slide, 7.25, 2.71, 7.65, 2.71, THEME.copper, 1.0)
    add_line(slide, 8.75, 2.71, 9.34, 2.71, THEME.copper, 1.0)
    add_line(slide, 10.03, 2.88, 10.55, 3.55, THEME.copper, 1.0)
    add_line(slide, 10.38, 3.88, 9.4, 4.22, THEME.copper, 1.0)
    add_line(slide, 7.65, 2.92, 8.62, 4.22, "C3452C", 1.4)
    add_text(slide, "small interface error", 6.28, 3.52, 1.62, 0.13, size=8, color="C3452C", font=THEME.mono, bold=True)
    add_rect(slide, 5.98, 5.32, 5.58, 0.48, fill=THEME.soft_copper, line=None)
    add_text(slide, "Fix: staged training + geometric consistency at model boundaries", 6.18, 5.49, 4.52, 0.12, size=9, color=THEME.ink, font=THEME.mono)
    for i, (value, label) in enumerate(s.metrics):
        add_value_label(slide, value, label, 0.82 + i * 1.65, 5.52)
    add_footer(slide, idx, total)


def render_loop(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide)
    add_kicker(slide, s.kicker)
    add_title(slide, s.title, size=23, width=10.8)
    add_text(slide, s.subtitle, M, 1.58, 6.2, 0.36, size=10, color=THEME.muted, font=THEME.sans)
    series_a = s.labels[0] if len(s.labels) >= 1 else "candidate"
    series_b = s.labels[1] if len(s.labels) >= 2 else "baseline"
    img_x, img_y, img_w, img_h = 0.82, 2.18, 6.45, 3.58
    add_rect(slide, img_x - 0.1, img_y - 0.1, img_w + 0.2, img_h + 0.2, fill="FFFFFF", line=THEME.hair)
    image_path = evidence_image_path(s)
    if image_path:
        add_image_path_fit(slide, image_path, img_x, img_y, img_w, img_h, border=False)
    else:
        # Native fallback chart for decks without a result figure.
        x0, y0, cw, ch = 1.28, 2.98, 4.85, 1.9
        add_line(slide, x0, y0 + ch, x0 + cw, y0 + ch, THEME.hair, 0.7)
        add_line(slide, x0, y0, x0, y0 + ch, THEME.hair, 0.7)
        pts1 = [(x0 + i * cw / 3, y0 + ch * (1 - v)) for i, v in enumerate((0.10, 0.43, 0.62, 0.76))]
        pts2 = [(x0 + i * cw / 3, y0 + ch * (1 - v)) for i, v in enumerate((0.18, 0.52, 0.46, 0.40))]
        for pts, color in ((pts1, "5F4BFF"), (pts2, "00A95C")):
            for a, b in zip(pts, pts[1:]):
                add_line(slide, a[0], a[1], b[0], b[1], color, 1.4)
    evidence = s.evidence
    if evidence:
        notes = add_evidence_pins(slide, evidence, img_x, img_y, img_w, img_h, dark=False)
        add_text(slide, "inspection notes", 7.65, 2.22, 2.2, 0.1, size=7, color=THEME.muted, font=THEME.mono, bold=True)
        add_note_ledger(slide, notes, 7.65, 2.52, 3.08, dark=False)
    add_text(slide, f"{series_a} vs {series_b}", 0.9, 5.98, 1.8, 0.12, size=8, color=THEME.copper, font=THEME.mono, bold=True)
    add_source_note(slide, evidence.source if evidence else "", 2.5, 5.98, 3.8)
    add_multiline(slide, s.bullets[:3], 7.65, 3.98, 4.15, 1.08, size=10, bullet=True)
    loop = (("probe", 7.65), ("mine", 8.7), ("train", 9.75), ("gate", 10.8))
    for i, (label, x) in enumerate(loop):
        add_node(slide, label, x, 5.44, 0.68, 0.28, mono=True, size=7)
        if i < 3:
            add_line(slide, x + 0.72, 5.58, x + 1.0, 5.58, THEME.copper, 0.9)
    add_rect(slide, 7.65, 6.02, 3.85, 0.38, fill=THEME.soft_copper, line=None)
    add_text(slide, "production failures become supervised signal", 7.84, 6.14, 3.45, 0.12, size=8, color=THEME.ink, font=THEME.mono)
    add_footer(slide, idx, total)


def render_product(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide, dark=True)
    add_grid(slide, dark=True)
    add_kicker(slide, s.kicker)
    add_text(slide, s.title, M, 0.72, 9.4, 0.44, size=24, color=THEME.paper, font=THEME.sans, bold=True)
    add_text(slide, s.subtitle, M, 1.42, 6.8, 0.36, size=10, color="C7D1D8", font=THEME.sans)
    img_x, img_y, img_w, img_h = 0.72, 2.04, 7.38, 3.98
    add_rect(slide, img_x - 0.12, img_y - 0.12, img_w + 0.24, img_h + 0.24, fill="202832", line="3E4B56")
    image_path = evidence_image_path(s)
    if image_path:
        add_image_path_fit(slide, image_path, img_x, img_y, img_w, img_h, border=False)
    evidence = s.evidence
    if evidence:
        add_text(slide, evidence.caption or "Product screenshot", img_x, 6.18, 5.0, 0.26, size=8, color=THEME.copper, font=THEME.mono, bold=True)
    else:
        add_text(slide, "Product screenshot", img_x, 6.22, 1.6, 0.12, size=8, color=THEME.copper, font=THEME.mono, bold=True)
    if evidence:
        add_source_note(slide, evidence.source, img_x, 6.58, 3.8)
        notes = add_evidence_pins(slide, evidence, img_x, img_y, img_w, img_h, dark=True)
        add_text(slide, "product signals", 8.58, 2.08, 2.2, 0.1, size=7, color="8A96A8", font=THEME.mono, bold=True)
        add_note_ledger(slide, notes, 8.58, 2.36, 2.95, dark=True)
    add_rect(slide, 8.52, 3.78, 3.75, 1.86, fill="202832", line="3E4B56")
    for i, bullet in enumerate(s.bullets[:3]):
        y = 4.08 + i * 0.42
        title = bullet.split(":")[0].strip()
        if i == 2:
            title = "Benchmark loop"
        add_text(slide, f"0{i+1}", 8.78, y, 0.28, 0.1, size=7, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, title[:24], 9.12, y - 0.02, 2.2, 0.12, size=8, color=THEME.paper, font=THEME.sans, bold=True)
    add_line(slide, 8.78, 5.76, 11.45, 5.76, "4A5967", 0.65)
    add_text(slide, "User-facing state is where benchmark gaps become concrete.", 8.78, 5.98, 3.2, 0.14, size=8, color=THEME.copper, font=THEME.mono, bold=True)
    add_footer(slide, idx, total)


def render_stack(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide)
    add_kicker(slide, s.kicker)
    add_title(slide, s.title, size=26, width=10.8)
    add_text(slide, s.subtitle, M, 1.86, 5.8, 0.45, size=12, color=THEME.muted, font=THEME.sans)
    for i, label in enumerate(s.labels):
        y = 2.64 + i * 0.62
        add_text(slide, f"{i+1:02d}", M, y, 0.35, 0.12, size=8, color=THEME.copper, font=THEME.mono, bold=True)
        add_text(slide, label, 1.05, y - 0.04, 3.2, 0.2, size=15, color=THEME.ink, font=THEME.serif)
        add_line(slide, 4.7, y + 0.07, 11.8, y + 0.07, THEME.hair, 0.55)
    add_rect(slide, 7.18, 2.38, 4.38, 2.02, fill=THEME.graphite, line=None)
    add_multiline(slide, s.bullets, 7.55, 2.78, 3.55, 1.18, size=11, color="DDE5EA", bullet=False)
    add_footer(slide, idx, total)


def render_close(slide, s: Slide, idx: int, total: int) -> None:
    canvas(slide, dark=True)
    add_grid(slide, dark=True)
    add_kicker(slide, s.kicker)
    title_is_long = len(s.title) > 58
    title_h = 0.92 if title_is_long else 0.58
    title_size = 27 if title_is_long else 31
    subtitle_y = 2.08 if title_is_long else 1.78
    add_text(slide, s.title, M, 0.86, 9.4, title_h, size=title_size, color=THEME.paper, font=THEME.serif)
    add_text(slide, s.subtitle, M, subtitle_y, 7.8, 0.48, size=13, color="C7D1D8", font=THEME.sans)
    for i, (label, bullet) in enumerate(zip(s.labels, s.bullets)):
        x = 0.85 + (i % 2) * 5.8
        y = 3.05 + (i // 2) * 1.15
        add_text(slide, label, x, y, 2.6, 0.22, size=17, color=THEME.paper, font=THEME.serif)
        add_text(slide, bullet, x, y + 0.38, 4.6, 0.34, size=10, color="AAB6BF", font=THEME.sans)
    add_line(slide, M, 6.16, W - M, 6.16, "53606B", 0.65)
    contact = s.note or ACTIVE_DECK.contact
    add_text(slide, f"{ACTIVE_DECK.title} | {contact}", M, 6.48, 7.6, 0.18, size=10, color=THEME.paper, font=THEME.mono)
    add_text(slide, f"{idx}/{total}", W - 0.98, 7.1, 0.42, 0.16, size=8, color="AAB6BF", font=THEME.sans, align=PP_ALIGN.RIGHT)


RENDERERS = {
    "cover": render_cover,
    "thesis": render_thesis,
    "map": render_map,
    "split": render_split,
    "process": render_process,
    "evidence": render_evidence,
    "matrix": render_matrix,
    "system": render_system,
    "loop": render_loop,
    "product": render_product,
    "stack": render_stack,
    "close": render_close,
}


def render_pptx(output_path: Path, deck: Deck | None = None) -> Path:
    global ACTIVE_DECK, ACTIVE_CACHE_DIR
    deck = deck or DEFAULT_DECK
    ACTIVE_DECK = deck
    ACTIVE_CACHE_DIR = output_path.parent / "asset-cache"
    prs = Presentation()
    prs.slide_width = emu(W)
    prs.slide_height = emu(H)
    blank = prs.slide_layouts[6]
    total = len(deck.slides)
    for idx, slide_data in enumerate(deck.slides, start=1):
        slide = prs.slides.add_slide(blank)
        renderer = RENDERERS[slide_data.kind]
        if slide_data.kind == "cover":
            renderer(slide, slide_data)
        else:
            renderer(slide, slide_data, idx, total)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
